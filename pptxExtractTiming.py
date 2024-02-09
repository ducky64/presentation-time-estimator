import argparse
from typing import *
from xml.etree.ElementTree import Element

from pptx import Presentation  # type: ignore
from pptx.oxml.xmlchemy import BaseOxmlElement  # type: ignore

parser = argparse.ArgumentParser(description='Extract script in comments and provide statistics including timing.')
parser.add_argument('pptx', type=str,
                    help='PPTX file to read')
args = parser.parse_args()


WPM_RATE = 120
BEAT_TIME = 1.0  # seconds

def words_time(word_count: int) -> float:  # in seconds
  return word_count / WPM_RATE * 60


def sec_to_str(sec: float) -> str:
  return f"{int(sec / 60)}:{int(sec % 60):02d}"


prs = Presentation(args.pptx)
slides = prs.slides

# https://stackoverflow.com/questions/45599473/python-pptx-how-to-access-presentation-section-titles-programmatically
# Because xmlchemy wasn't built with any extensibility or the possibility of nonstandard namespaces
# this gets monkey patched in place
# https://github.com/scanny/python-pptx/blob/master/pptx/oxml/xmlchemy.py
nsmap = {
  'p14': "http://schemas.microsoft.com/office/powerpoint/2010/main"
}
section_lists = super(BaseOxmlElement, prs._element).xpath('.//p14:sectionLst', namespaces=nsmap)
assert len(section_lists) == 1
section_list = cast(Element, section_lists[0])

all_str = ""
summary_str = ""

cumulative_time = 0.0  # in seconds
cumulative_words = 0
cumulative_breaks = 0.0  # in seconds
for i, section in enumerate(section_list):
  name = section.attrib['name']
  section_slides = section.xpath('.//p14:sldId', namespaces=nsmap)
  section_ids = [int(slide.attrib['id']) for slide in section_slides]

  section_str = ""
  section_narration = ""
  section_words = 0
  section_breaks = 0
  for slide_id in section_ids:
    slide = prs.slides.get(slide_id)

    if slide._element.attrib.get('show', '1') == '0':
      continue

    if slide.shapes.title is not None:
      title = slide.shapes.title.text
    else:
      title = ""
    section_str += f"### {prs.slides.index(slide)} {title}\n"

    notes = slide.notes_slide.notes_text_frame.text
    # split by line so we recognize special commands
    notes_lines = notes.replace('\r', ' ').split('\n')
    notes_words = 0
    notes_breaks = 0
    for line in notes_lines:
      line = line.strip()
      if line.startswith('#pause: '):
        unit = line[8:]
        time = {
          'beat': BEAT_TIME
        }
        notes_breaks += time[unit]
      else:
        notes_words += len(list(filter(None, line.split())))

    section_str += f"_{notes_words} words ({sec_to_str(words_time(notes_words))}) + {sec_to_str(notes_breaks)} breaks_\n"
    section_str += "\n"
    section_str += notes
    section_str += "\n"

    section_words += notes_words
    section_breaks += notes_breaks

  section_words_time = words_time(section_words)

  all_str += f"## <u>{name}</u>\n"
  all_str += f"_{sec_to_str(words_time(cumulative_words))} + {section_words} words ({sec_to_str(section_words_time)}) + {sec_to_str(section_breaks)} breaks_\n"
  all_str += "\n"
  all_str += "\n"
  all_str += section_str
  all_str += "\n"
  all_str += "\n"

  summary_str += f"**{name}**: {sec_to_str(cumulative_time)} + {section_words} words ({sec_to_str(section_words_time)}) + {sec_to_str(section_breaks)} breaks\n"

  cumulative_time += section_words_time + section_breaks
  cumulative_words += section_words
  cumulative_breaks += section_breaks


summary_str += '\n'
summary_str += f"**TOTAL**: {sec_to_str(cumulative_time)}: {cumulative_words} words ({sec_to_str(words_time(cumulative_words))}) + {sec_to_str(cumulative_breaks)} breaks\n"

all_str = summary_str + all_str

print(summary_str)

with open(args.pptx + "_timing.md", 'w', encoding='utf-8') as outfile:
  outfile.write(all_str)
